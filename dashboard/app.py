"""
Hormuz Shock Simulator — Streamlit Dashboard
"""
from __future__ import annotations

import json
import sys
from io import BytesIO
from pathlib import Path

import pandas as pd
import plotly.graph_objects as go
import streamlit as st

# ---------------------------------------------------------------------------
# Path setup — ensure project root is importable
# ---------------------------------------------------------------------------
_ROOT = Path(__file__).resolve().parents[1]
if str(_ROOT) not in sys.path:
    sys.path.insert(0, str(_ROOT))

from data.baseline.loader import load_baseline
from engine.propagation import PropagationEngine
from engine.sentiment import SentimentCascade
from engine.shock_params import Severity, ShockScenario
from geo.build_layers import build_all_layers
from geo.kepler_config import build_kepler_config

# ---------------------------------------------------------------------------
# Page config
# ---------------------------------------------------------------------------
st.set_page_config(
    page_title="Hormuz Shock Simulator",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ---------------------------------------------------------------------------
# Cached loaders
# ---------------------------------------------------------------------------
@st.cache_data(show_spinner="Loading BEDA baseline data…")
def get_baseline() -> dict:
    return load_baseline()


@st.cache_data(show_spinner="Running propagation model…")
def run_propagation(duration: int, severity: str, name: str) -> tuple[pd.DataFrame, pd.Series]:
    scenario = ShockScenario(
        duration_weeks=duration,
        severity=Severity(severity),
        scenario_name=name,
    )
    baseline = get_baseline()
    engine = PropagationEngine(scenario, baseline)
    return engine.propagate(), engine.oil_price_path()


@st.cache_data(show_spinner="Running sentiment cascade…")
def run_sentiment(duration: int, severity: str, name: str) -> pd.DataFrame:
    scenario = ShockScenario(
        duration_weeks=duration,
        severity=Severity(severity),
        scenario_name=name,
    )
    return SentimentCascade(scenario).propagate()


# ---------------------------------------------------------------------------
# Sidebar — scenario controls
# ---------------------------------------------------------------------------
st.sidebar.title("Scenario Parameters")

with st.sidebar:
    scenario_name = st.text_input("Scenario Name", value="Hormuz Closure — Base Case")

    severity_map = {
        "PARTIAL (×1.25)": "PARTIAL",
        "MODERATE (×1.55)": "MODERATE",
        "SEVERE (×1.85)": "SEVERE",
        "EXTREME (×2.20)": "EXTREME",
    }
    severity_label = st.selectbox(
        "Severity",
        list(severity_map.keys()),
        index=2,
    )
    severity = severity_map[severity_label]

    duration = st.slider("Duration (weeks)", min_value=1, max_value=52, value=16, step=1)

    snapshot_week = st.select_slider(
        "Map Snapshot Week",
        options=[4, 13, 26, 52],
        value=13,
    )

    st.divider()
    st.caption("**Scenario Comparison** — add up to 3 scenarios")

    compare_scenarios: list[dict] = []
    for i in range(3):
        with st.expander(f"Comparison Scenario {i + 1}", expanded=(i == 0)):
            cn = st.text_input(f"Name {i}", value=f"Scenario {i + 1}", key=f"cn_{i}")
            cs = st.selectbox(f"Severity {i}", list(severity_map.keys()), index=i, key=f"cs_{i}")
            cd = st.slider(f"Duration {i} (wks)", 1, 52, [4, 8, 16][i], key=f"cd_{i}")
            compare_scenarios.append({"name": cn, "severity": severity_map[cs], "duration": cd})


# ---------------------------------------------------------------------------
# Run model for primary scenario
# ---------------------------------------------------------------------------
ts_df, oil_path = run_propagation(duration, severity, scenario_name)
sentiment_df = run_sentiment(duration, severity, scenario_name)

scenario = ShockScenario(
    duration_weeks=duration,
    severity=Severity(severity),
    scenario_name=scenario_name,
)

# Build geo layers
baseline = get_baseline()
sa4_geojson, arcs_geojson, freight_geojson = build_all_layers(
    ts_df, scenario.get_oil_multiplier(), baseline
)
kepler_cfg = build_kepler_config(snapshot_week)

# ---------------------------------------------------------------------------
# Tabs
# ---------------------------------------------------------------------------
tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs(
    ["Impact Map", "Indicator Impacts", "Sentiment Cascade", "Scenario Comparison", "Sources", "How to Use"]
)

# ---------------------------------------------------------------------------
# TAB 1: pydeck map (deck.gl — same engine as Kepler.gl)
# ---------------------------------------------------------------------------
with tab1:
    st.subheader(f"Brisbane Economic Impact — Week {snapshot_week} Snapshot")
    st.caption(
        f"Scenario: **{scenario_name}** | Severity: **{severity}** | "
        f"Duration: **{duration} weeks** | Oil multiplier: **\u00d7{scenario.get_oil_multiplier():.2f}**"
    )

    try:
        import pydeck as pdk

        impact_field = f"impact_score_wk{snapshot_week}"

        # ---- Layer 1: SA4 choropleth (PolygonLayer, sequential purple) ----
        # Absolute scale: 25 = full dark purple (25% avg deviation from baseline).
        # This means colours change meaningfully with severity, duration, and snapshot week.
        SCALE_MAX = 25.0

        sa4_poly_data = []
        for feat in sa4_geojson["features"]:
            score = feat["properties"].get(impact_field, 0)
            t = min(score / SCALE_MAX, 1.0)
            # Purple gradient: light (#f2f0f7) → dark (#54278f)
            r = int(242 - t * (242 - 84))
            g = int(240 - t * (240 - 39))
            b = int(247 - t * (247 - 143))
            # PolygonLayer expects a flat list of [lon, lat] coordinate pairs
            coords = feat["geometry"]["coordinates"][0]
            sa4_poly_data.append({
                "polygon": coords,
                "name": feat["properties"]["sa4_name"],
                "fill_color": [r, g, b, 180],
                "tooltip_html": f"<b>{feat['properties']['sa4_name']}</b><br/>Impact index: {score:.1f} / 100<br/><span style='font-size:11px;color:#aaa'>avg % deviation from baseline, exposure-weighted</span>",
            })

        geojson_layer = pdk.Layer(
            "PolygonLayer",
            data=sa4_poly_data,
            get_polygon="polygon",
            get_fill_color="fill_color",
            get_line_color=[255, 255, 255, 100],
            line_width_min_pixels=1,
            pickable=True,
            auto_highlight=True,
        )

        # ---- Layer 2: Route arcs — one pdk.Layer per disruption level ----
        # pydeck ArcLayer doesn't reliably honour list-valued field accessors for
        # get_source_color, so we split by level and use a constant colour per layer.
        DISRUPTION_COLOR = {
            "LOW":      [33, 150, 243, 200],
            "MODERATE": [255, 152, 0,  220],
            "HIGH":     [244, 67, 54,  230],
            "CRITICAL": [183, 28, 28,  255],
        }

        # Build arc records grouped by disruption level
        arc_by_level: dict[str, list] = {lvl: [] for lvl in DISRUPTION_COLOR}
        for feat in arcs_geojson["features"]:
            coords = feat["geometry"]["coordinates"]
            props = feat["properties"]
            level = props["disruption_level"]
            arc_by_level.setdefault(level, []).append({
                "source_position": coords[0],
                "target_position": coords[1],
                "arc_weight": props["arc_weight"],
                "tooltip_html": (
                    f"<b>BNE → {props['destination_name']}</b><br/>"
                    f"Disruption: <b>{level}</b>"
                    + (f" <span style='color:#aaa;font-size:11px'>(base: {props['base_disruption_level']})</span>"
                       if level != props['base_disruption_level'] else "")
                    + f"<br/>Route weight: {props['arc_weight']:.2f}<br/><br/>"
                    f"<span style='font-size:11px;max-width:280px;display:block'>"
                    f"{props['reasoning']}</span>"
                ),
            })

        arc_layers = [
            pdk.Layer(
                "ArcLayer",
                data=records,
                get_source_position="source_position",
                get_target_position="target_position",
                get_source_color=color,
                get_target_color=color,
                get_width="arc_weight * 2",
                pickable=True,
                auto_highlight=True,
            )
            for level, color in DISRUPTION_COLOR.items()
            if (records := arc_by_level.get(level, []))
        ]

        # ---- Layer 3: Freight corridors (PathLayer, orange heat) ----
        FUEL_COLOR = {
            "LOW":      [255, 236, 210, 180],
            "MEDIUM":   [255, 179, 71,  200],
            "HIGH":     [255, 140, 0,   220],
            "CRITICAL": [178, 34, 34,   255],
        }
        freight_path_data = []
        for feat in freight_geojson["features"]:
            props = feat["properties"]
            freight_path_data.append({
                "path": feat["geometry"]["coordinates"],
                "color": FUEL_COLOR.get(props["fuel_intensity"], [255, 140, 0, 200]),
                "tooltip_html": (
                    f"<b>{props['name']}</b><br/>"
                    f"Fuel intensity: {props['fuel_intensity']}<br/>"
                    f"Freight cost increase: {props['cost_impact_pct']:.1f}%"
                ),
            })

        path_layer = pdk.Layer(
            "PathLayer",
            data=freight_path_data,
            get_path="path",
            get_color="color",
            get_width=4,
            width_min_pixels=3,
            pickable=True,
        )

        view_state = pdk.ViewState(
            longitude=153.02,
            latitude=-27.47,
            zoom=9.5,
            pitch=0,
            bearing=0,
        )

        deck = pdk.Deck(
            layers=[geojson_layer, *arc_layers, path_layer],
            initial_view_state=view_state,
            tooltip={
                "html": "{tooltip_html}",
                "style": {"backgroundColor": "#1a1a2e", "color": "white", "fontSize": "13px", "padding": "8px", "lineHeight": "1.6"},
            },
            map_style="https://basemaps.cartocdn.com/gl/dark-matter-gl-style/style.json",
        )
        # Debug strip — shows live disruption levels so you can confirm data changes
        with st.expander("Route disruption levels (current scenario)", expanded=False):
            for feat in arcs_geojson["features"]:
                p = feat["properties"]
                st.write(f"BNE → {p['destination_name']}: **{p['disruption_level']}** (base: {p['base_disruption_level']})")

        # Render as raw HTML to bypass Streamlit's component diffing,
        # which silently skips re-renders when deck data changes.
        deck_html = deck.to_html(as_string=True, notebook_display=False)
        st.components.v1.html(deck_html, height=620)

        # Legend
        col_l1, col_l2, col_l3 = st.columns(3)
        with col_l1:
            st.markdown("**SA4 Choropleth** — light purple = low impact, dark purple = ≥25% avg deviation from baseline")
        with col_l2:
            st.markdown("**Route Arcs** — blue=low, amber=moderate, red=high/critical disruption")
        with col_l3:
            st.markdown("**Freight Corridors** — orange intensity = fuel cost impact %")

        st.divider()
        with st.expander("Route disruption rationale — hover an arc or expand here"):
            from geo.build_layers import ROUTE_DESTINATIONS
            LEVEL_COLOR = {
                "LOW": "#2196F3", "MODERATE": "#FF9800",
                "HIGH": "#F44336", "CRITICAL": "#B71C1C",
            }
            for dest in ROUTE_DESTINATIONS:
                color = LEVEL_COLOR.get(dest["disruption_level"], "#888")
                st.markdown(
                    f"**BNE → {dest['name']}** &nbsp;"
                    f"<span style='background:{color};color:white;padding:2px 8px;"
                    f"border-radius:4px;font-size:12px'>{dest['disruption_level']}</span>",
                    unsafe_allow_html=True,
                )
                st.caption(dest["reasoning"])

    except Exception as e:
        st.warning(f"pydeck rendering issue: {e}. Showing plotly fallback.")
        fig = go.Figure()
        for feat in sa4_geojson["features"]:
            coords = feat["geometry"]["coordinates"][0]
            lons = [c[0] for c in coords] + [coords[0][0]]
            lats = [c[1] for c in coords] + [coords[0][1]]
            score = feat["properties"].get(f"impact_score_wk{snapshot_week}", 0)
            fig.add_trace(go.Scattergeo(
                lon=lons, lat=lats, mode="lines+text",
                text=[feat["properties"]["sa4_name"]] + [""] * (len(lons) - 1),
                name=feat["properties"]["sa4_name"],
                line_color=f"rgba(108,20,138,{min(score * 2, 1.0):.2f})",
                textposition="top center",
                showlegend=True,
            ))
        for feat in arcs_geojson["features"]:
            coords = feat["geometry"]["coordinates"]
            fig.add_trace(go.Scattergeo(
                lon=[coords[0][0], coords[1][0]],
                lat=[coords[0][1], coords[1][1]],
                mode="lines",
                name=feat["properties"]["destination"],
                line={"width": feat["properties"]["arc_weight"] * 3, "color": "#FF8C00"},
                showlegend=False,
            ))
        fig.update_layout(
            geo={
                "scope": "world",
                "center": {"lon": 153.02, "lat": -27.47},
                "projection_scale": 15,
                "showland": True, "landcolor": "#2a2a2a",
                "bgcolor": "#1a1a1a", "showocean": True, "oceancolor": "#1e2a3a",
            },
            paper_bgcolor="#1a1a1a",
            height=600,
            margin={"l": 0, "r": 0, "t": 0, "b": 0},
        )
        st.plotly_chart(fig, width='stretch')

# ---------------------------------------------------------------------------
# TAB 2: Indicator impact table
# ---------------------------------------------------------------------------
with tab2:
    st.subheader("Indicator Impact — % Change from Baseline")
    st.caption("Values show % change from BEDA baseline. Red = deterioration, green = improvement.")

    # Build flat baseline lookup: indicator → (current_value, unit_hint)
    flat_baseline_vals: dict[str, float] = {}
    for cat, inds in baseline.items():
        for ind, vals in inds.items():
            cv = vals.get("current_value")
            if isinstance(cv, (int, float)) and cv != 0:
                flat_baseline_vals[ind] = float(cv)

    display_weeks = [w for w in [4, 8, 13, 26, 52] if w <= 52]

    # Convert raw deltas → % change from baseline
    pct_df = ts_df.loc[display_weeks].copy()
    for col in pct_df.columns:
        bv = flat_baseline_vals.get(col)
        if bv and abs(bv) > 0:
            pct_df[col] = (pct_df[col] / abs(bv)) * 100
        else:
            pct_df[col] = 0.0

    snapshot = pct_df.T.copy()
    snapshot.columns = [f"Wk {w}" for w in display_weeks]
    snapshot.index.name = "Indicator"
    snapshot = snapshot.reset_index()

    # Add category and baseline value columns
    cat_map: dict[str, str] = {}
    for cat, inds in baseline.items():
        for ind in inds:
            cat_map[ind] = cat
    snapshot.insert(0, "Category", snapshot["Indicator"].map(cat_map).fillna("Other"))
    snapshot.insert(2, "Baseline", snapshot["Indicator"].map(flat_baseline_vals))

    categories = sorted(snapshot["Category"].unique())
    selected_cats = st.multiselect(
        "Filter by Category", categories, default=categories[:4]
    )
    filtered = snapshot[snapshot["Category"].isin(selected_cats)] if selected_cats else snapshot

    week_cols = [c for c in filtered.columns if c.startswith("Wk")]

    def _color_pct(val):
        if not isinstance(val, (int, float)):
            return ""
        if val < -20:
            return "background-color: #8B0000; color: white"
        if val < -5:
            return "background-color: #CC4444; color: white"
        if val < 0:
            return "background-color: #FFCCCC"
        if val > 20:
            return "background-color: #006400; color: white"
        if val > 5:
            return "background-color: #44AA44; color: white"
        if val > 0:
            return "background-color: #CCFFCC"
        return ""

    fmt = {c: "{:+.1f}%" for c in week_cols}
    fmt["Baseline"] = "{:,.1f}"
    styled = (
        filtered.set_index(["Category", "Indicator"])
        .style.map(_color_pct, subset=week_cols)
        .format(fmt)
    )

    st.dataframe(styled, width='stretch', height=500)

    # Oil path chart
    st.subheader("Oil Price Multiplier Path")
    fig_oil = go.Figure()
    fig_oil.add_trace(go.Scatter(
        x=list(oil_path.index), y=list(oil_path.values),
        mode="lines+markers", name="Oil Price Multiplier",
        line={"color": "#FF8C00", "width": 2},
        fill="tozeroy", fillcolor="rgba(255,140,0,0.2)",
    ))
    fig_oil.add_hline(y=1.0, line_dash="dash", line_color="gray", annotation_text="Baseline")
    fig_oil.update_layout(
        xaxis_title="Week", yaxis_title="Multiplier (1.0 = baseline)",
        paper_bgcolor="#0e1117", plot_bgcolor="#1a1a1a",
        font_color="white", height=300,
    )
    st.plotly_chart(fig_oil, width='stretch')

# ---------------------------------------------------------------------------
# TAB 3: Sentiment cascade
# ---------------------------------------------------------------------------
with tab3:
    st.subheader("Sentiment Cascade Timeline")
    st.caption(
        "Chain: Oil Shock → Pump Price → Disposable Income → Consumer Sentiment "
        "→ Retail Spend → Business Confidence → Hiring → Investment Pipeline"
    )

    colors = [
        "#FF4500", "#FF8C00", "#FFD700", "#90EE90",
        "#00CED1", "#4169E1", "#9370DB", "#FF69B4",
    ]

    fig_sent = go.Figure()
    for i, col in enumerate(sentiment_df.columns):
        fig_sent.add_trace(go.Scatter(
            x=list(sentiment_df.index),
            y=list(sentiment_df[col]),
            mode="lines",
            name=col,
            line={"color": colors[i % len(colors)], "width": 2},
        ))

    fig_sent.add_hline(y=0, line_dash="dot", line_color="gray", annotation_text="No change")
    fig_sent.update_layout(
        xaxis_title="Week",
        yaxis_title="Delta from baseline (index points)",
        paper_bgcolor="#0e1117",
        plot_bgcolor="#1a1a1a",
        font_color="white",
        height=480,
        legend={"orientation": "v", "x": 1.01, "y": 1},
    )
    st.plotly_chart(fig_sent, width='stretch')

    # Summary table
    from engine.sentiment import SentimentCascade as SC
    stats = SC(scenario).summary_stats()
    st.dataframe(stats.style.format({"peak_abs_delta": "{:.2f}", "trough_value": "{:.2f}"}),
                 width='stretch')

# ---------------------------------------------------------------------------
# TAB 4: Scenario comparison
# ---------------------------------------------------------------------------
with tab4:
    st.subheader("Scenario Comparison")

    all_scenarios = [
        {"name": scenario_name, "severity": severity, "duration": duration}
    ] + compare_scenarios[:3]

    # Select indicator for radar
    all_inds = list(ts_df.columns)
    radar_inds = st.multiselect(
        "Select indicators for radar chart (8–12 recommended)",
        all_inds,
        default=[
            "Jet fuel price - Singapore Kerosene (USD/bbl)",
            "Brisbane CBD hotel occupancy rate (%)",
            "Queensland retail trade turnover (monthly, $B)",
            "Unemployment rate - Brisbane SA4 (combined, %)",
            "Brisbane median house price - CoreLogic (AUD)",
            "NAB Business Confidence Index - Queensland",
            "RACQ 91 ULP Brisbane - annual average (cpl)",
            "Brisbane CBD RevPAR (AUD)",
        ],
    )

    radar_week = st.select_slider("Radar at Week", options=[4, 13, 26, 52], value=13, key="radar_wk")

    if radar_inds:
        fig_radar = go.Figure()
        scenario_table_rows = []

        for sc_def in all_scenarios:
            sc_ts, _ = run_propagation(sc_def["duration"], sc_def["severity"], sc_def["name"])
            week_idx = min(radar_week, len(sc_ts))
            vals = sc_ts.loc[week_idx, [i for i in radar_inds if i in sc_ts.columns]].tolist()
            # Normalise to pct of baseline for radar
            baseline_vals = []
            for ind in [i for i in radar_inds if i in sc_ts.columns]:
                bv = 0.0
                for cat_data in baseline.values():
                    if ind in cat_data:
                        v = cat_data[ind]["current_value"]
                        bv = float(v) if isinstance(v, (int, float)) else 1.0
                        break
                baseline_vals.append(abs(bv) if bv != 0 else 1.0)
            pct_vals = [abs(v / bv) * 100 for v, bv in zip(vals, baseline_vals)]

            fig_radar.add_trace(go.Scatterpolar(
                r=pct_vals + pct_vals[:1],
                theta=[i for i in radar_inds if i in sc_ts.columns] +
                      [[i for i in radar_inds if i in sc_ts.columns][0]],
                fill="toself",
                name=sc_def["name"],
            ))

            scenario_table_rows.append({
                "Scenario": sc_def["name"],
                "Severity": sc_def["severity"],
                "Duration (wks)": sc_def["duration"],
                "Oil Mult": ShockScenario(sc_def["duration"], Severity(sc_def["severity"])).get_oil_multiplier(),
                **{ind[:30]: round(sc_ts.loc[week_idx, ind], 3)
                   for ind in radar_inds if ind in sc_ts.columns},
            })

        fig_radar.update_layout(
            polar={
                "bgcolor": "#1a1a1a",
                "radialaxis": {"visible": True, "range": [0, max(50, 1)]},
            },
            paper_bgcolor="#0e1117",
            font_color="white",
            height=500,
            showlegend=True,
        )
        st.plotly_chart(fig_radar, width='stretch')

        comparison_df = pd.DataFrame(scenario_table_rows)
        st.dataframe(comparison_df.set_index("Scenario"), width='stretch')

    # -----------------------------------------------------------------------
    # Export button
    # -----------------------------------------------------------------------
    st.divider()
    st.subheader("Export Results")

    col_xl, col_pptx = st.columns(2)

    with col_xl:
        if st.button("Generate Excel Export"):
            xl_buf = BytesIO()
            with pd.ExcelWriter(xl_buf, engine="openpyxl") as writer:
                ts_df.to_excel(writer, sheet_name="Propagation Deltas")
                sentiment_df.to_excel(writer, sheet_name="Sentiment Cascade")
                oil_path.to_frame().to_excel(writer, sheet_name="Oil Price Path")
                # Summary sheet
                snap = ts_df.loc[[4, 13, 26, 52]].T
                snap.columns = ["Week 4", "Week 13", "Week 26", "Week 52"]
                snap.to_excel(writer, sheet_name="Snapshot Summary")
            xl_buf.seek(0)
            st.download_button(
                "Download Excel",
                data=xl_buf,
                file_name=f"hormuz_shock_{scenario_name.replace(' ', '_')}.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )

    with col_pptx:
        if st.button("Generate PPTX Summary"):
            try:
                from pptx import Presentation
                from pptx.util import Inches, Pt, Emu
                from pptx.dml.color import RGBColor
                from pptx.enum.text import PP_ALIGN
                import io

                prs = Presentation()
                prs.slide_width = Inches(13.33)
                prs.slide_height = Inches(7.5)

                blank_layout = prs.slide_layouts[6]  # blank

                def add_text_box(slide, text, left, top, width, height, font_size=18,
                                 bold=False, color=(255, 255, 255), bg_color=None):
                    txBox = slide.shapes.add_textbox(
                        Inches(left), Inches(top), Inches(width), Inches(height)
                    )
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = text
                    p.font.size = Pt(font_size)
                    p.font.bold = bold
                    p.font.color.rgb = RGBColor(*color)
                    if bg_color:
                        fill = txBox.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor(*bg_color)
                    return txBox

                # Slide 1: Title
                slide1 = prs.slides.add_slide(blank_layout)
                bg = slide1.background
                bg.fill.solid()
                bg.fill.fore_color.rgb = RGBColor(15, 17, 23)
                add_text_box(slide1, "Hormuz Shock Simulator", 0.5, 1.5, 12, 1.5,
                             font_size=36, bold=True, color=(255, 140, 0))
                add_text_box(slide1, f"Scenario: {scenario_name}", 0.5, 3.0, 12, 0.8,
                             font_size=24, color=(200, 200, 200))
                add_text_box(
                    slide1,
                    f"Severity: {severity}  |  Duration: {duration} weeks  |  "
                    f"Oil multiplier: ×{scenario.get_oil_multiplier():.2f}",
                    0.5, 3.9, 12, 0.7, font_size=18, color=(180, 180, 180)
                )
                add_text_box(slide1, "Brisbane Economic Development Agency (BEDA)", 0.5, 6.5, 12, 0.5,
                             font_size=12, color=(120, 120, 120))

                # Slide 2: Key indicators at week 13
                slide2 = prs.slides.add_slide(blank_layout)
                slide2.background.fill.solid()
                slide2.background.fill.fore_color.rgb = RGBColor(15, 17, 23)
                add_text_box(slide2, "Key Indicator Impacts — Week 13", 0.5, 0.3, 12, 0.8,
                             font_size=28, bold=True, color=(255, 140, 0))

                key_inds = [
                    "Jet fuel price - Singapore Kerosene (USD/bbl)",
                    "RACQ 91 ULP Brisbane - annual average (cpl)",
                    "Brisbane CBD hotel occupancy rate (%)",
                    "Queensland retail trade turnover (monthly, $B)",
                    "Unemployment rate - Brisbane SA4 (combined, %)",
                    "NAB Business Confidence Index - Queensland",
                    "Brisbane median house price - CoreLogic (AUD)",
                    "Westpac-Melbourne Institute Consumer Sentiment Index",
                ]
                for idx, ind in enumerate(key_inds):
                    if ind in ts_df.columns:
                        col_x = 0.5 if idx % 2 == 0 else 6.8
                        row_y = 1.3 + (idx // 2) * 1.2
                        val = ts_df.loc[13, ind]
                        color = (255, 100, 100) if val < 0 else (100, 220, 100)
                        add_text_box(slide2, f"{ind[:38]}", col_x, row_y, 6.0, 0.5,
                                     font_size=11, color=(180, 180, 180))
                        add_text_box(slide2, f"{val:+.3f}", col_x, row_y + 0.45, 6.0, 0.45,
                                     font_size=14, bold=True, color=color)

                # Slide 3: Sentiment cascade summary
                slide3 = prs.slides.add_slide(blank_layout)
                slide3.background.fill.solid()
                slide3.background.fill.fore_color.rgb = RGBColor(15, 17, 23)
                add_text_box(slide3, "Sentiment Cascade — Peak Impacts", 0.5, 0.3, 12, 0.8,
                             font_size=28, bold=True, color=(255, 140, 0))

                from engine.sentiment import SentimentCascade as SC
                stats = SC(scenario).summary_stats()
                for idx, (node, row) in enumerate(stats.iterrows()):
                    y_pos = 1.3 + idx * 0.68
                    add_text_box(slide3, f"{node}", 0.5, y_pos, 6.5, 0.4,
                                 font_size=12, color=(200, 200, 200))
                    add_text_box(slide3,
                                 f"Peak Δ: {row['peak_abs_delta']:.2f}  at Week {int(row['peak_week'])}",
                                 7.2, y_pos, 5.5, 0.4, font_size=12, color=(255, 200, 100))

                # Slide 4: Scenario comparison table
                slide4 = prs.slides.add_slide(blank_layout)
                slide4.background.fill.solid()
                slide4.background.fill.fore_color.rgb = RGBColor(15, 17, 23)
                add_text_box(slide4, "Scenario Comparison (Week 13 deltas)", 0.5, 0.3, 12, 0.8,
                             font_size=28, bold=True, color=(255, 140, 0))

                sc_summary_rows = []
                for sc_def in all_scenarios:
                    sc_ts2, _ = run_propagation(sc_def["duration"], sc_def["severity"], sc_def["name"])
                    row_data = {"Scenario": sc_def["name"][:25], "Sev": sc_def["severity"],
                                "Dur": sc_def["duration"]}
                    for ind in key_inds[:4]:
                        if ind in sc_ts2.columns:
                            row_data[ind[:18]] = f"{sc_ts2.loc[13, ind]:+.2f}"
                    sc_summary_rows.append(row_data)

                for ridx, row_data in enumerate(sc_summary_rows):
                    y_pos = 1.3 + ridx * 1.2
                    row_text = "  |  ".join(f"{k}: {v}" for k, v in row_data.items())
                    add_text_box(slide4, row_text, 0.5, y_pos, 12.3, 0.9,
                                 font_size=11, color=(200, 200, 200))

                pptx_buf = BytesIO()
                prs.save(pptx_buf)
                pptx_buf.seek(0)
                st.download_button(
                    "Download PPTX",
                    data=pptx_buf,
                    file_name=f"hormuz_shock_{scenario_name.replace(' ', '_')}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            except Exception as ex:
                st.error(f"PPTX generation failed: {ex}")

# ---------------------------------------------------------------------------
# TAB 5: Sources
# ---------------------------------------------------------------------------
with tab5:
    st.subheader("Data Sources and Model Calibration")
    st.caption("All indicators, transfer coefficients, and shock calibrations used in this tool.")

    st.markdown("### Baseline Data")
    st.markdown("""
| Source | What it provides |
|---|---|
| **BEDA_Enriched.xlsx** | Current values, prior values, and YoY change for all 13 indicator categories. Primary input for baseline calibration. |
| **ABS ASGS 2021** | SA4 boundary polygons for Brisbane choropleth map. Fetched via ABS ArcGIS REST API and cached locally. |
| **ABS Labour Force Survey** | Unemployment and underemployment rates for Brisbane SA4 and Queensland. |
| **CoreLogic** | Brisbane median house price and unit price. |
| **RACQ Fuel Price Watch** | 91 ULP and diesel Brisbane annual averages (cpl). |
| **ABS CPI** | Brisbane all-groups and food subgroup CPI, year-on-year. |
| **Westpac–Melbourne Institute** | Consumer Sentiment Index, used as a leading indicator for retail spend. |
| **NAB Business Survey** | Business Confidence Index — Queensland. |
| **Brisbane Airport Corporation** | Total passengers, international load factor, weekly seat capacity. |
| **ABS Overseas Arrivals and Departures** | International visitor arrivals Queensland. |
| **ABS Retail Trade** | Queensland monthly retail trade turnover. |
| **AEMO** | Queensland wholesale electricity spot price (Q4 average). |
| **ABS Building Approvals** | Residential and new dwelling approvals — Queensland and Brisbane. |
""")

    st.markdown("### Shock Calibration")
    st.markdown("""
Transfer coefficients in `engine/propagation.py` are calibrated from the following historical events:

| Reference event | Oil price move | Key Brisbane outcomes used |
|---|---|---|
| **Russia–Ukraine 2022** | Brent +143% peak, settled +63% YoY | Brisbane 91 ULP peak 231.6 cpl; CPI food +7.6% YoY Q3 2022; NAB Business Confidence −20pp within 3 months; Consumer Sentiment fell to 83.7 by Aug 2022 |
| **Libya supply loss 2011** | +55% over 3 months | Moderate scenario calibration anchor |
| **Gulf tanker attacks 2019** | +25% spike, ~2 weeks | Partial scenario calibration anchor |
| **Hormuz closure literature** | Modelled +120% (20% global supply removed) | Extreme scenario ceiling |

The **Brisbane car-dependency multiplier (1.18)** reflects ~80% private vehicle trip share vs ~65% national average, applied to all fuel-linked indicators.

Jet fuel transfer coefficient (0.95) is calibrated to Singapore Kerosene spot, which directly prices BNE airline fuel contracts.
""")

    st.markdown("### Route Arc Disruption Levels")
    st.markdown("""
Base disruption levels are calibrated to a **Moderate severity** scenario (oil ×1.55) and scaled up or down
by one tier per severity step. The base levels reflect each hub's structural dependency on Hormuz-linked supply chains:

| Route | Base level | Primary reason |
|---|---|---|
| BNE–SIN | CRITICAL | Singapore refines ~1.5M bbl/day of Hormuz crude; prices BNE jet fuel contracts |
| BNE–NRT | HIGH | Japan imports ~90% of oil via Hormuz |
| BNE–ICN | HIGH | South Korea sources ~70% of crude via Hormuz |
| BNE–HKG | HIGH | Cathay Pacific exposed; HKG is key cargo re-export hub |
| BNE–PVG | MODERATE | China has ~90 days strategic reserves + Russia pipeline alternative |
| BNE–LAX | LOW | US sources <10% of crude from Gulf; Pacific route unaffected |
| BNE–AKL | LOW | Short-haul trans-Tasman; route operates regardless of Hormuz |
""")

    st.markdown("### Freight Corridor Cost Multipliers")
    st.markdown("""
Freight cost impact is computed as `(oil_multiplier − 1.0) × fuel_intensity_factor`.
Fuel intensity factors: LOW = 0.5, MEDIUM = 0.75, HIGH = 1.0, CRITICAL = 1.25.

Port of Brisbane access is rated CRITICAL due to diesel-heavy container logistics.
Bruce Highway and Pacific Motorway are rated HIGH. Freight rail to Rocklea is MEDIUM.
""")

# ---------------------------------------------------------------------------
# TAB 6: How to Use
# ---------------------------------------------------------------------------
with tab6:
    st.subheader("How to Use This Tool")
    st.markdown("""
This tool simulates the economic impact on Brisbane if the Strait of Hormuz closes,
propagating an oil price shock through BEDA's indicator set and modelling the human
sentiment cascade that follows. All outputs show **deltas from the BEDA baseline** —
what changes, not what the absolute level becomes.
""")

    st.markdown("### Sidebar Parameters")
    st.markdown("""
| Parameter | What it controls |
|---|---|
| **Scenario Name** | Label for the scenario — used in exports and comparison tab. |
| **Severity** | How severe the oil price shock is. Drives the oil price multiplier (PARTIAL ×1.25 through EXTREME ×2.20) and determines route arc disruption colours. |
| **Duration (weeks)** | How long the Strait stays closed before recovery begins. Affects how far lagged indicators (unemployment, housing, investment) have time to move — most visible at snapshot weeks 26 and 52. |
| **Map Snapshot Week** | Which week to display on the SA4 choropleth. Week 4 shows the acute fuel/CPI shock; week 52 shows compounded effects after all lags have played out. |
| **Comparison Scenarios** | Define up to 3 additional scenarios for side-by-side comparison on the Scenario Comparison tab. |
""")

    st.markdown("### Tabs")

    st.markdown("#### Impact Map")
    st.markdown("""
Three overlapping layers:
- **SA4 choropleth** — Brisbane's 9 statistical sub-regions shaded by composite impact index (0–100).
  Colour scale is absolute: full dark purple = 25% average deviation from baseline across all indicators.
  Hover a region to see its impact score.
- **Route arcs** — Flight/trade routes from Brisbane Airport (BNE) to 7 international hubs.
  Colour reflects disruption severity: blue = LOW, amber = MODERATE, red = HIGH, dark red = CRITICAL.
  Levels shift with scenario severity. Hover an arc to read the full disruption rationale.
- **Freight corridors** — Brisbane's major road and rail freight arteries.
  Orange intensity = estimated freight cost increase percentage. Hover to see the corridor name and cost impact.

Expand **Route disruption rationale** below the map to read all 7 arc explanations at once.
""")

    st.markdown("#### Indicator Impacts")
    st.markdown("""
A table showing **% change from BEDA baseline** for every indicator at weeks 4, 8, 13, 26, and 52.
The **Baseline** column shows the current value so you can interpret the scale.

Colour thresholds: dark red = >20% deterioration, light red = 5–20%, light green = 5–20% improvement,
dark green = >20% improvement. Use the **Filter by Category** selector to focus on specific indicator groups.

The **Oil Price Multiplier Path** chart below the table shows how the oil price evolves over 52 weeks
for the current scenario — the ramp-up, plateau, and exponential decay back toward baseline.
""")

    st.markdown("#### Sentiment Cascade")
    st.markdown("""
Models the human behavioural chain that follows an oil shock:

**Oil shock → Pump price → Disposable income → Consumer sentiment →
Retail spend → Business confidence → Hiring → Investment pipeline**

Each node has its own lag and transfer coefficient. The chart shows how each node's index value
deviates from baseline over 52 weeks. Nodes further along the chain peak later but can persist
longer than the oil shock itself. The summary table below the chart shows peak deviation and
the week it occurs for each node.
""")

    st.markdown("#### Scenario Comparison")
    st.markdown("""
Compare up to 4 scenarios side by side (the primary scenario plus up to 3 defined in the sidebar).
Select indicators for the **radar chart** — 8 to 12 works best. Use the **Radar at Week** slider
to compare how scenarios diverge at different points in time.

The comparison table below the radar shows raw indicator deltas at the selected week for each scenario.

Use the **Export** section to download results as Excel (all time series) or a 4-slide PPTX summary.
""")

    st.markdown("### Impact Score Methodology")
    st.markdown("""
The SA4 impact index is a **weighted average of % changes from baseline** across 7 indicator categories:
Aviation & Connectivity, Domestic Travel, Visitor Economy, Consumer Cost of Living,
Retail & Consumer Spending, Labour Market, and Housing & Real Estate.

Each category's % change is weighted by that region's **exposure weight** — for example,
Brisbane - West has a higher weight for Consumer Cost of Living (1.05) due to high car dependency,
while Brisbane Inner City has a higher weight for Visitor Economy (1.00) due to CBD hotel and
events concentration.

A score of 15 means indicators in that region are deviating by an average of 15% from baseline,
after adjusting for regional exposure. The scale is absolute (not relative between regions),
so comparing scores across severity levels and snapshot weeks is meaningful.
""")
